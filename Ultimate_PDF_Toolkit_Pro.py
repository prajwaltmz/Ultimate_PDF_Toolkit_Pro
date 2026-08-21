import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog
import os
import sys

# Make the UI crisp and "super HD" on high DPI displays (Windows only)
try:
    import ctypes
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except Exception:
    pass

try:
    import fitz  # PyMuPDF
    from pypdf import PdfWriter, PdfReader
    from PIL import Image, ImageTk
    from pdf2docx import Converter
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    messagebox.showerror("Dependency Error", f"Missing dependency: {e}\nPlease run 'pip install -r requirements.txt'")
    sys.exit(1)

# --- REUSABLE PDF VIEWER WINDOW ---
class PDFViewerWindow(tk.Toplevel):
    def __init__(self, parent, pdf_path):
        super().__init__(parent)
        self.title(f"Viewing: {os.path.basename(pdf_path)}")
        self.geometry("800x900")
        self.configure(bg="#2b2b2b")
        self.pdf_path = pdf_path
        
        try:
            self.doc = fitz.open(pdf_path)
            self.current_page = 0
        except Exception as e:
            messagebox.showerror("Error", f"Failed to open PDF:\n{e}")
            self.destroy()
            return
            
        # Controls
        ctrl_frame = tk.Frame(self, bg="#2b2b2b")
        ctrl_frame.pack(fill=tk.X, pady=10)
        
        ttk.Button(ctrl_frame, text="⬅️ Previous", command=self.prev_page).pack(side=tk.LEFT, padx=20)
        self.lbl_page = tk.Label(ctrl_frame, text=f"Page {self.current_page + 1} / {len(self.doc)}", bg="#2b2b2b", fg="white", font=("Segoe UI", 12, "bold"))
        self.lbl_page.pack(side=tk.LEFT, expand=True)
        ttk.Button(ctrl_frame, text="Next ➡️", command=self.next_page).pack(side=tk.RIGHT, padx=20)
        
        # Canvas for image
        self.canvas = tk.Canvas(self, bg="#3c3f41", highlightthickness=0)
        self.canvas.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        self.bind("<Configure>", self.on_resize)
        self.show_page()

    def show_page(self):
        if not hasattr(self, 'doc'): return
        page = self.doc.load_page(self.current_page)
        # High resolution render
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        mode = "RGBA" if pix.alpha else "RGB"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        
        # Calculate resize to fit canvas
        cw = self.canvas.winfo_width()
        ch = self.canvas.winfo_height()
        if cw < 100 or ch < 100:
            cw, ch = 760, 800
            
        img.thumbnail((cw, ch), Image.Resampling.LANCZOS)
        self.tk_img = ImageTk.PhotoImage(img)
        
        self.canvas.delete("all")
        self.canvas.create_image(cw//2, ch//2, image=self.tk_img, anchor=tk.CENTER)
        self.lbl_page.config(text=f"Page {self.current_page + 1} / {len(self.doc)}")

    def next_page(self):
        if self.current_page < len(self.doc) - 1:
            self.current_page += 1
            self.show_page()

    def prev_page(self):
        if self.current_page > 0:
            self.current_page -= 1
            self.show_page()
            
    def on_resize(self, event):
        if event.widget == self:
            self.show_page()

# --- MAIN APPLICATION ---
class UltimatePDFToolkit:
    def __init__(self, root):
        self.root = root
        self.root.title("Ultimate PDF Toolkit Pro")
        self.root.geometry("850x650")
        self.root.configure(bg="#1e1e1e")
        
        # Style
        style = ttk.Style()
        style.theme_use('clam')
        style.configure("TNotebook", background="#1e1e1e", borderwidth=0)
        style.configure("TNotebook.Tab", background="#2d2d2d", foreground="#aaaaaa", padding=[10, 5], font=('Segoe UI', 10, 'bold'))
        style.map("TNotebook.Tab", background=[("selected", "#007acc")], foreground=[("selected", "white")])
        style.configure("TFrame", background="#1e1e1e")
        style.configure("TButton", background="#007acc", foreground="white", font=('Segoe UI', 10, 'bold'), padding=5)
        style.map("TButton", background=[("active", "#005999")])
        
        # Title
        tk.Label(root, text="Ultimate PDF Toolkit Pro", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 24, "bold")).pack(pady=10)
        
        self.notebook = ttk.Notebook(root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.init_organize_tab()
        self.init_merge_pdf_tab()
        self.init_merge_word_tab()
        self.init_convert_tab()
        self.init_image_tab()
        self.init_pptx_tab()
        self.init_split_tab()
        self.init_lock_tab()
        self.init_annotate_tab()

    def create_preview_btn(self, parent, path_getter):
        def cmd():
            p = path_getter()
            if p and os.path.exists(p) and p.lower().endswith(".pdf"):
                PDFViewerWindow(self.root, p)
            else:
                messagebox.showerror("Error", "Please select a valid PDF file first.")
        ttk.Button(parent, text="👁️ Preview PDF", command=cmd).pack(pady=5)

    # 1. ORGANIZE PDF
    def init_organize_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="1. Organize")
        tk.Label(f, text="Add a PDF, view pages, remove or reorder them.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.org_path = tk.StringVar()
        ttk.Button(f, text="Select PDF", command=self.load_org_pdf).pack(pady=5)
        self.create_preview_btn(f, lambda: self.org_path.get())
        
        self.org_list = tk.Listbox(f, selectmode=tk.SINGLE, bg="#2b2b2b", fg="white", font=("Segoe UI", 10), height=10)
        self.org_list.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        bf = tk.Frame(f, bg="#1e1e1e")
        bf.pack(pady=5)
        ttk.Button(bf, text="Move Up ⬆️", command=lambda: self.move_org(-1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Move Down ⬇️", command=lambda: self.move_org(1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Delete Page 🗑️", command=self.del_org).pack(side=tk.LEFT, padx=5)
        ttk.Button(f, text="Save Organized PDF 💾", command=self.save_org).pack(pady=10)

    def load_org_pdf(self):
        p = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")])
        if p:
            self.org_path.set(p)
            self.org_list.delete(0, tk.END)
            reader = PdfReader(p)
            for i in range(len(reader.pages)):
                self.org_list.insert(tk.END, f"Page {i+1}")

    def move_org(self, delta):
        sel = self.org_list.curselection()
        if not sel: return
        idx = sel[0]
        new_idx = idx + delta
        if 0 <= new_idx < self.org_list.size():
            item = self.org_list.get(idx)
            self.org_list.delete(idx)
            self.org_list.insert(new_idx, item)
            self.org_list.selection_set(new_idx)

    def del_org(self):
        sel = self.org_list.curselection()
        if sel:
            self.org_list.delete(sel[0])

    def save_org(self):
        if not self.org_path.get() or self.org_list.size() == 0:
            messagebox.showerror("Error", "No pages to save!")
            return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if out:
            reader = PdfReader(self.org_path.get())
            writer = PdfWriter()
            for item in self.org_list.get(0, tk.END):
                p_num = int(item.split()[1]) - 1
                writer.add_page(reader.pages[p_num])
            with open(out, "wb") as f:
                writer.write(f)
            messagebox.showinfo("Success", "Organized PDF saved successfully!")

    # 2. MERGE PDF
    def init_merge_pdf_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="2. Merge PDF")
        tk.Label(f, text="Merge multiple PDF files into one.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.pdf_list = tk.Listbox(f, selectmode=tk.SINGLE, bg="#2b2b2b", fg="white", font=("Segoe UI", 10))
        self.pdf_list.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        bf = tk.Frame(f, bg="#1e1e1e")
        bf.pack(pady=5)
        ttk.Button(bf, text="Add PDFs", command=self.add_pdf_files).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Move Up", command=lambda: self.move_list_item(self.pdf_list, -1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Move Down", command=lambda: self.move_list_item(self.pdf_list, 1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Remove", command=lambda: self.pdf_list.delete(self.pdf_list.curselection()) if self.pdf_list.curselection() else None).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Clear", command=lambda: self.pdf_list.delete(0, tk.END)).pack(side=tk.LEFT, padx=5)
        
        self.create_preview_btn(f, lambda: self.pdf_list.get(self.pdf_list.curselection()) if self.pdf_list.curselection() else None)
        ttk.Button(f, text="Merge PDFs 💾", command=self.merge_pdfs).pack(pady=10)

    def add_pdf_files(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF files", "*.pdf")])
        for file in files:
            self.pdf_list.insert(tk.END, file)

    def merge_pdfs(self):
        files = self.pdf_list.get(0, tk.END)
        if not files:
            messagebox.showerror("Error", "No PDF files selected!")
            return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if out:
            writer = PdfWriter()
            for file in files:
                writer.append(file)
            with open(out, "wb") as f:
                writer.write(f)
            messagebox.showinfo("Success", "PDFs merged successfully!")

    # 3. MERGE WORD
    def init_merge_word_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="3. Merge Word")
        tk.Label(f, text="Merge multiple Word documents (.docx, .doc) via MS Word COM.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.word_list = tk.Listbox(f, selectmode=tk.SINGLE, bg="#2b2b2b", fg="white", font=("Segoe UI", 10))
        self.word_list.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        bf = tk.Frame(f, bg="#1e1e1e")
        bf.pack(pady=5)
        ttk.Button(bf, text="Add Word Docs", command=self.add_word_files).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Move Up", command=lambda: self.move_list_item(self.word_list, -1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Move Down", command=lambda: self.move_list_item(self.word_list, 1)).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Remove", command=lambda: self.word_list.delete(self.word_list.curselection()) if self.word_list.curselection() else None).pack(side=tk.LEFT, padx=5)
        ttk.Button(bf, text="Clear", command=lambda: self.word_list.delete(0, tk.END)).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(f, text="Merge Word Docs 💾", command=self.merge_word_docs).pack(pady=10)

    def add_word_files(self):
        files = filedialog.askopenfilenames(filetypes=[("Word documents", "*.docx;*.doc")])
        for file in files:
            self.word_list.insert(tk.END, file)

    def merge_word_docs(self):
        files = self.word_list.get(0, tk.END)
        if not files:
            messagebox.showerror("Error", "No Word files selected!")
            return
        out = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word Document", "*.docx")])
        if out:
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                main_doc = word.Documents.Open(os.path.abspath(files[0]))
                for file in files[1:]:
                    word.Selection.EndKey(Unit=6)  # wdStory
                    word.Selection.InsertBreak(Type=7)  # wdPageBreak
                    word.Selection.InsertFile(FileName=os.path.abspath(file))
                main_doc.SaveAs(os.path.abspath(out))
                main_doc.Close()
                word.Quit()
                messagebox.showinfo("Success", "Word documents merged successfully!")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to merge Word documents:\n{e}\n(Ensure MS Word is installed)")

    # 4. CONVERT PDF <-> WORD
    def init_convert_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="4. PDF ↔ Word")
        
        tk.Label(f, text="Convert between PDF and Word", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        # PDF to Word
        p2w = tk.LabelFrame(f, text=" PDF to Word ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        p2w.pack(fill=tk.X, padx=20, pady=10)
        self.p2w_file = tk.StringVar()
        tk.Entry(p2w, textvariable=self.p2w_file, width=50, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(p2w, text="Browse", command=lambda: self.p2w_file.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(side=tk.LEFT, padx=5)
        ttk.Button(p2w, text="Convert to Word", command=self.convert_pdf_to_word).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(p2w, lambda: self.p2w_file.get())
        
        # Word to PDF
        w2p = tk.LabelFrame(f, text=" Word to PDF ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        w2p.pack(fill=tk.X, padx=20, pady=10)
        self.w2p_file = tk.StringVar()
        tk.Entry(w2p, textvariable=self.w2p_file, width=50, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(w2p, text="Browse", command=lambda: self.w2p_file.set(filedialog.askopenfilename(filetypes=[("Word", "*.docx;*.doc")]))).pack(side=tk.LEFT, padx=5)
        ttk.Button(w2p, text="Convert to PDF", command=self.convert_word_to_pdf).pack(side=tk.LEFT, padx=5)

    def convert_pdf_to_word(self):
        pdf = self.p2w_file.get()
        if not pdf: return
        out = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word", "*.docx")])
        if out:
            try:
                cv = Converter(pdf)
                cv.convert(out, start=0, end=None)
                cv.close()
                messagebox.showinfo("Success", "Converted PDF to Word successfully!")
            except Exception as e:
                messagebox.showerror("Error", f"Conversion failed:\n{e}")

    def convert_word_to_pdf(self):
        doc_file = self.w2p_file.get()
        if not doc_file: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if out:
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(doc_file))
                doc.SaveAs(os.path.abspath(out), FileFormat=17)  # 17 = wdFormatPDF
                doc.Close()
                word.Quit()
                messagebox.showinfo("Success", "Converted Word to PDF successfully!")
            except Exception as e:
                messagebox.showerror("Error", f"Conversion failed:\n{e}\n(Ensure MS Word is installed)")

    # 5. IMAGE <-> PDF
    def init_image_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="5. Images ↔ PDF")
        tk.Label(f, text="Convert Images to PDF or Extract Images from PDF", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        # Images to PDF
        i2p = tk.LabelFrame(f, text=" Images to PDF ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        i2p.pack(fill=tk.X, padx=20, pady=10)
        self.img_list = []
        self.lbl_img_count = tk.Label(i2p, text="0 images selected", bg="#1e1e1e", fg="white")
        self.lbl_img_count.pack(side=tk.LEFT, padx=5)
        ttk.Button(i2p, text="Select Images", command=self.select_images).pack(side=tk.LEFT, padx=5)
        ttk.Button(i2p, text="Create PDF", command=self.images_to_pdf).pack(side=tk.LEFT, padx=5)
        
        # PDF to Images
        p2i = tk.LabelFrame(f, text=" PDF to Images (Extract Pages) ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        p2i.pack(fill=tk.X, padx=20, pady=10)
        self.p2i_file = tk.StringVar()
        tk.Entry(p2i, textvariable=self.p2i_file, width=40, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(p2i, text="Browse PDF", command=lambda: self.p2i_file.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(p2i, lambda: self.p2i_file.get())
        ttk.Button(p2i, text="Extract All Pages to Images", command=self.pdf_to_images).pack(side=tk.LEFT, padx=5)

    def select_images(self):
        self.img_list = filedialog.askopenfilenames(filetypes=[("Image files", "*.jpg;*.jpeg;*.png;*.bmp")])
        self.lbl_img_count.config(text=f"{len(self.img_list)} images selected")

    def images_to_pdf(self):
        if not self.img_list: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if out:
            imgs = [Image.open(i).convert('RGB') for i in self.img_list]
            imgs[0].save(out, save_all=True, append_images=imgs[1:])
            messagebox.showinfo("Success", "Images converted to PDF successfully!")

    def pdf_to_images(self):
        pdf = self.p2i_file.get()
        if not pdf: return
        out_dir = filedialog.askdirectory()
        if out_dir:
            doc = fitz.open(pdf)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                pix.save(os.path.join(out_dir, f"page_{i+1}.png"))
            messagebox.showinfo("Success", f"All pages exported to images in:\n{out_dir}")

    # 6. PDF TO PPTX
    def init_pptx_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="6. PDF to PPTX")
        tk.Label(f, text="Convert PDF pages directly into PowerPoint (.pptx) slides.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.pptx_pdf = tk.StringVar()
        tk.Entry(f, textvariable=self.pptx_pdf, width=50, bg="#2b2b2b", fg="white").pack(pady=5)
        ttk.Button(f, text="Browse PDF", command=lambda: self.pptx_pdf.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(pady=5)
        self.create_preview_btn(f, lambda: self.pptx_pdf.get())
        ttk.Button(f, text="Convert to PowerPoint 📊", command=self.convert_pdf_to_pptx).pack(pady=10)

    def convert_pdf_to_pptx(self):
        pdf = self.pptx_pdf.get()
        if not pdf: return
        out = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if out:
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            doc = fitz.open(pdf)
            
            temp_files = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                t_img = f"temp_p_{i}.png"
                pix.save(t_img)
                temp_files.append(t_img)
                
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(t_img, Inches(0.5), Inches(0.5), width=Inches(9))
                
            prs.save(out)
            for t in temp_files:
                if os.path.exists(t): os.remove(t)
            messagebox.showinfo("Success", "PDF converted to PowerPoint presentation successfully!")

    # 7. SPLIT PDF
    def init_split_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="7. Split PDF")
        tk.Label(f, text="Extract specific page ranges or split into single pages.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.split_pdf = tk.StringVar()
        tk.Entry(f, textvariable=self.split_pdf, width=50, bg="#2b2b2b", fg="white").pack(pady=5)
        ttk.Button(f, text="Browse PDF", command=lambda: self.split_pdf.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(pady=5)
        self.create_preview_btn(f, lambda: self.split_pdf.get())
        
        rf = tk.Frame(f, bg="#1e1e1e")
        rf.pack(pady=10)
        tk.Label(rf, text="Extract Range (e.g. 1-5):", bg="#1e1e1e", fg="white").pack(side=tk.LEFT, padx=5)
        self.range_var = tk.StringVar()
        tk.Entry(rf, textvariable=self.range_var, width=15, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(rf, text="Extract Range", command=self.extract_range).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(f, text="Split Entire Document into Single Pages", command=self.split_all).pack(pady=10)

    def extract_range(self):
        pdf = self.split_pdf.get()
        rng = self.range_var.get()
        if not pdf or not rng: return
        try:
            parts = rng.split('-')
            start = int(parts[0]) - 1
            end = int(parts[1]) if len(parts) > 1 else start + 1
            
            reader = PdfReader(pdf)
            writer = PdfWriter()
            for i in range(start, min(end, len(reader.pages))):
                writer.add_page(reader.pages[i])
                
            out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
            if out:
                with open(out, "wb") as f:
                    writer.write(f)
                messagebox.showinfo("Success", "Extracted pages saved successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to extract range: {e}")

    def split_all(self):
        pdf = self.split_pdf.get()
        if not pdf: return
        out_dir = filedialog.askdirectory()
        if out_dir:
            reader = PdfReader(pdf)
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                with open(os.path.join(out_dir, f"page_{i+1}.pdf"), "wb") as f:
                    writer.write(f)
            messagebox.showinfo("Success", f"All pages saved as separate PDFs in:\n{out_dir}")

    # 8. LOCK & UNLOCK
    def init_lock_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="8. Lock / Unlock")
        tk.Label(f, text="Add password protection or remove passwords from PDF.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        # Lock
        lf = tk.LabelFrame(f, text=" Lock (Encrypt PDF) ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        lf.pack(fill=tk.X, padx=20, pady=10)
        self.lock_pdf = tk.StringVar()
        tk.Entry(lf, textvariable=self.lock_pdf, width=40, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(lf, text="Browse", command=lambda: self.lock_pdf.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(side=tk.LEFT, padx=5)
        ttk.Button(lf, text="Set Password & Lock 🔒", command=self.do_lock).pack(side=tk.LEFT, padx=5)
        
        # Unlock
        uf = tk.LabelFrame(f, text=" Unlock (Remove Password) ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        uf.pack(fill=tk.X, padx=20, pady=10)
        self.unlock_pdf = tk.StringVar()
        tk.Entry(uf, textvariable=self.unlock_pdf, width=40, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(uf, text="Browse", command=lambda: self.unlock_pdf.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(side=tk.LEFT, padx=5)
        ttk.Button(uf, text="Unlock & Save Unencrypted 🔓", command=self.do_unlock).pack(side=tk.LEFT, padx=5)

    def do_lock(self):
        pdf = self.lock_pdf.get()
        if not pdf: return
        pwd = simpledialog.askstring("Password", "Enter password to lock PDF:", show='*')
        if pwd:
            out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
            if out:
                reader = PdfReader(pdf)
                writer = PdfWriter()
                for p in reader.pages: writer.add_page(p)
                writer.encrypt(pwd)
                with open(out, "wb") as f: writer.write(f)
                messagebox.showinfo("Success", "PDF encrypted and locked successfully!")

    def do_unlock(self):
        pdf = self.unlock_pdf.get()
        if not pdf: return
        pwd = simpledialog.askstring("Password", "Enter current password for PDF:", show='*')
        if pwd is not None:
            try:
                reader = PdfReader(pdf)
                if reader.is_encrypted:
                    reader.decrypt(pwd)
                writer = PdfWriter()
                for p in reader.pages: writer.add_page(p)
                out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
                if out:
                    with open(out, "wb") as f: writer.write(f)
                    messagebox.showinfo("Success", "PDF unlocked and saved successfully!")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to unlock PDF:\n{e}")

    # 9. ANNOTATE & EXTRACT
    def init_annotate_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="9. Annotate & Extract")
        tk.Label(f, text="Extract text or search & highlight keywords in PDF.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.ann_pdf = tk.StringVar()
        tk.Entry(f, textvariable=self.ann_pdf, width=50, bg="#2b2b2b", fg="white").pack(pady=5)
        ttk.Button(f, text="Browse PDF", command=lambda: self.ann_pdf.set(filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")]))).pack(pady=5)
        self.create_preview_btn(f, lambda: self.ann_pdf.get())
        
        bf = tk.Frame(f, bg="#1e1e1e")
        bf.pack(pady=10)
        ttk.Button(bf, text="Extract All Text to .txt", command=self.extract_text).pack(side=tk.LEFT, padx=10)
        
        hf = tk.LabelFrame(f, text=" Search & Highlight Keyword ", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 10, "bold"), padx=10, pady=10)
        hf.pack(fill=tk.X, padx=20, pady=10)
        self.kw_var = tk.StringVar()
        tk.Label(hf, text="Keyword:", bg="#1e1e1e", fg="white").pack(side=tk.LEFT, padx=5)
        tk.Entry(hf, textvariable=self.kw_var, width=20, bg="#2b2b2b", fg="white").pack(side=tk.LEFT, padx=5)
        ttk.Button(hf, text="Highlight & Save 🖍️", command=self.highlight_keyword).pack(side=tk.LEFT, padx=5)

    def extract_text(self):
        pdf = self.ann_pdf.get()
        if not pdf: return
        out = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text File", "*.txt")])
        if out:
            doc = fitz.open(pdf)
            full_text = ""
            for i, page in enumerate(doc):
                full_text += f"\n--- Page {i+1} ---\n" + page.get_text()
            with open(out, "w", encoding="utf-8") as f:
                f.write(full_text)
            messagebox.showinfo("Success", "Extracted text saved to file!")

    def highlight_keyword(self):
        pdf = self.ann_pdf.get()
        kw = self.kw_var.get().strip()
        if not pdf or not kw: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if out:
            doc = fitz.open(pdf)
            found = 0
            for page in doc:
                inst = page.search_for(kw)
                for rect in inst:
                    page.add_highlight_annot(rect)
                    found += 1
            doc.save(out)
            messagebox.showinfo("Success", f"Highlighted {found} instances of '{kw}' and saved to new PDF!")

    def move_list_item(self, listbox, delta):
        sel = listbox.curselection()
        if not sel: return
        idx = sel[0]
        new_idx = idx + delta
        if 0 <= new_idx < listbox.size():
            item = listbox.get(idx)
            listbox.delete(idx)
            listbox.insert(new_idx, item)
            listbox.selection_set(new_idx)

if __name__ == "__main__":
    root = tk.Tk()
    app = UltimatePDFToolkit(root)
    root.mainloop()
