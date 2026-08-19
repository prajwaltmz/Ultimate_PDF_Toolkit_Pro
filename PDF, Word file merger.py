import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog
import os
import sys

# Make the UI crisp and "super HD" on high DPI displays (Windows only)
try:
    import ctypes
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except Exception:
    pass

try:
    import fitz  # PyMuPDF
    from pypdf import PdfWriter, PdfReader
    from PIL import Image, ImageTk
    from pdf2docx import Converter
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    messagebox.showerror("Dependency Error", f"Missing dependency: {e}\nPlease run 'pip install pymupdf pypdf Pillow pdf2docx python-pptx pywin32'")
    sys.exit(1)

# --- REUSABLE PDF VIEWER WINDOW ---
class PDFViewerWindow(tk.Toplevel):
    def __init__(self, parent, pdf_path):
        super().__init__(parent)
        self.title(f"Viewing: {os.path.basename(pdf_path)}")
        self.geometry("800x900")
        self.configure(bg="#2b2b2b")
        self.pdf_path = pdf_path
        
        try:
            self.doc = fitz.open(pdf_path)
            self.current_page = 0
        except Exception as e:
            messagebox.showerror("Error", f"Failed to open PDF:\n{e}")
            self.destroy()
            return
            
        # Controls
        ctrl_frame = tk.Frame(self, bg="#2b2b2b")
        ctrl_frame.pack(fill=tk.X, pady=10)
        
        ttk.Button(ctrl_frame, text="⬅️ Previous", command=self.prev_page).pack(side=tk.LEFT, padx=20)
        self.lbl_page = tk.Label(ctrl_frame, text=f"Page {self.current_page + 1} / {len(self.doc)}", bg="#2b2b2b", fg="white", font=("Segoe UI", 12, "bold"))
        self.lbl_page.pack(side=tk.LEFT, expand=True)
        ttk.Button(ctrl_frame, text="Next ➡️", command=self.next_page).pack(side=tk.RIGHT, padx=20)
        
        # Canvas for image
        self.canvas = tk.Canvas(self, bg="#3c3f41", highlightthickness=0)
        self.canvas.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        self.bind("<Configure>", self.on_resize)
        self.show_page()

    def show_page(self):
        if not hasattr(self, 'doc'): return
        page = self.doc.load_page(self.current_page)
        # High resolution render
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        mode = "RGBA" if pix.alpha else "RGB"
        img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
        
        # Calculate resize to fit canvas
        cw = self.canvas.winfo_width()
        ch = self.canvas.winfo_height()
        if cw < 100 or ch < 100:
            cw, ch = 760, 800
            
        img.thumbnail((cw, ch), Image.Resampling.LANCZOS)
        self.tk_img = ImageTk.PhotoImage(img)
        
        self.canvas.delete("all")
        self.canvas.create_image(cw//2, ch//2, image=self.tk_img, anchor=tk.CENTER)
        self.lbl_page.config(text=f"Page {self.current_page + 1} / {len(self.doc)}")

    def next_page(self):
        if self.current_page < len(self.doc) - 1:
            self.current_page += 1
            self.show_page()

    def prev_page(self):
        if self.current_page > 0:
            self.current_page -= 1
            self.show_page()
            
    def on_resize(self, event):
        if event.widget == self:
            self.show_page()

# --- MAIN APPLICATION ---
class UltimatePDFToolkit:
    def __init__(self, root):
        self.root = root
        self.root.title("Ultimate PDF Toolkit Pro V2 🚀")
        self.root.geometry("850x650")
        self.root.configure(bg="#1e1e1e")
        
        # Style
        style = ttk.Style()
        style.theme_use('clam')
        style.configure("TNotebook", background="#1e1e1e", borderwidth=0)
        style.configure("TNotebook.Tab", background="#2d2d2d", foreground="#aaaaaa", padding=[10, 5], font=('Segoe UI', 10, 'bold'))
        style.map("TNotebook.Tab", background=[("selected", "#007acc")], foreground=[("selected", "white")])
        style.configure("TFrame", background="#1e1e1e")
        style.configure("TButton", background="#007acc", foreground="white", font=('Segoe UI', 10, 'bold'), padding=5)
        style.map("TButton", background=[("active", "#005999")])
        
        # Title
        tk.Label(root, text="PDF Toolkit Pro V2", bg="#1e1e1e", fg="#007acc", font=("Segoe UI", 24, "bold")).pack(pady=10)
        
        self.notebook = ttk.Notebook(root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.init_organize_tab()
        self.init_merge_pdf_tab()
        self.init_merge_word_tab()
        self.init_convert_tab()
        self.init_image_tab()
        self.init_pptx_tab()
        self.init_split_tab()
        self.init_lock_tab()
        self.init_annotate_tab()

    def create_preview_btn(self, parent, path_getter):
        def cmd():
            p = path_getter()
            if p and os.path.exists(p) and p.lower().endswith(".pdf"):
                PDFViewerWindow(self.root, p)
            else:
                messagebox.showerror("Error", "Please select a valid PDF file first.")
        ttk.Button(parent, text="👁️ Preview PDF", command=cmd).pack(pady=5)

    # 1. ORGANIZE PDF
    def init_organize_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="1. Organize")
        tk.Label(f, text="Add a PDF, view pages, remove or reorder them.", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=10)
        
        self.org_pdf = None
        self.org_lbl = tk.Label(f, text="No PDF Selected", bg="#1e1e1e", fg="#007acc")
        self.org_lbl.pack()
        
        btn_f = tk.Frame(f, bg="#1e1e1e")
        btn_f.pack(pady=5)
        ttk.Button(btn_f, text="📁 Select PDF", command=self.org_select).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f, lambda: self.org_pdf)
        
        self.org_list = tk.Listbox(f, bg="#2d2d2d", fg="white", font=("Segoe UI", 11), selectbackground="#007acc")
        self.org_list.pack(fill=tk.BOTH, expand=True, padx=40, pady=10)
        
        ctrl_f = tk.Frame(f, bg="#1e1e1e")
        ctrl_f.pack(pady=5)
        ttk.Button(ctrl_f, text="⬆️ Move Up", command=lambda: self.list_move_up(self.org_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl_f, text="⬇️ Move Down", command=lambda: self.list_move_down(self.org_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl_f, text="❌ Remove Page", command=lambda: self.list_remove(self.org_list)).pack(side=tk.LEFT, padx=5)
        
        tk.Button(f, text="SAVE ORGANIZED PDF", bg="#28a745", fg="white", font=("Segoe UI", 12, "bold"), command=self.org_save).pack(fill=tk.X, padx=40, pady=10)

    def org_select(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.org_pdf = p
            self.org_lbl.config(text=os.path.basename(p))
            self.org_list.delete(0, tk.END)
            try:
                doc = fitz.open(p)
                for i in range(len(doc)):
                    self.org_list.insert(tk.END, f"Page {i+1}")
            except Exception as e:
                messagebox.showerror("Error", str(e))

    def org_save(self):
        if not self.org_pdf or self.org_list.size() == 0: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        
        try:
            reader = PdfReader(self.org_pdf)
            writer = PdfWriter()
            # extract original page index from "Page X"
            for i in range(self.org_list.size()):
                text = self.org_list.get(i)
                idx = int(text.split(" ")[1]) - 1
                writer.add_page(reader.pages[idx])
            with open(out, "wb") as f:
                writer.write(f)
            messagebox.showinfo("Success", "Organized PDF saved!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def list_move_up(self, lb):
        sel = lb.curselection()
        if sel and sel[0] > 0:
            idx = sel[0]
            val = lb.get(idx)
            lb.delete(idx)
            lb.insert(idx - 1, val)
            lb.selection_set(idx - 1)
            
    def list_move_down(self, lb):
        sel = lb.curselection()
        if sel and sel[0] < lb.size() - 1:
            idx = sel[0]
            val = lb.get(idx)
            lb.delete(idx)
            lb.insert(idx + 1, val)
            lb.selection_set(idx + 1)
            
    def list_remove(self, lb):
        sel = lb.curselection()
        if sel: lb.delete(sel[0])

    # 2. MERGE PDF
    def init_merge_pdf_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="2. Merge PDF")
        self.mpdf_list = tk.Listbox(f, bg="#2d2d2d", fg="white", font=("Segoe UI", 11), selectbackground="#007acc")
        self.mpdf_list.pack(fill=tk.BOTH, expand=True, padx=40, pady=20)
        
        ctrl = tk.Frame(f, bg="#1e1e1e")
        ctrl.pack(pady=5)
        ttk.Button(ctrl, text="➕ Add PDFs", command=lambda: self.add_to_list(self.mpdf_list, [("PDF", "*.pdf")])).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="⬆️ Up", command=lambda: self.list_move_up(self.mpdf_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="⬇️ Down", command=lambda: self.list_move_down(self.mpdf_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="❌ Remove", command=lambda: self.list_remove(self.mpdf_list)).pack(side=tk.LEFT, padx=5)
        
        tk.Button(f, text="MERGE PDFs", bg="#007acc", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_merge_pdf).pack(fill=tk.X, padx=40, pady=10)

    def add_to_list(self, lb, ftypes):
        for p in filedialog.askopenfilenames(filetypes=ftypes): lb.insert(tk.END, p)

    def do_merge_pdf(self):
        if self.mpdf_list.size() == 0: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        try:
            writer = PdfWriter()
            for i in range(self.mpdf_list.size()):
                writer.append(self.mpdf_list.get(i))
            writer.write(out)
            writer.close()
            messagebox.showinfo("Success", "PDFs merged successfully!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # 3. MERGE WORD
    def init_merge_word_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="3. Merge Word")
        tk.Label(f, text="Merge multiple Word files into ONE Word file.", bg="#1e1e1e", fg="white").pack(pady=10)
        
        self.mword_list = tk.Listbox(f, bg="#2d2d2d", fg="white", font=("Segoe UI", 11), selectbackground="#007acc")
        self.mword_list.pack(fill=tk.BOTH, expand=True, padx=40, pady=10)
        
        ctrl = tk.Frame(f, bg="#1e1e1e")
        ctrl.pack(pady=5)
        ttk.Button(ctrl, text="➕ Add Word Files", command=lambda: self.add_to_list(self.mword_list, [("Word", "*.docx *.doc")])).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="⬆️ Up", command=lambda: self.list_move_up(self.mword_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="⬇️ Down", command=lambda: self.list_move_down(self.mword_list)).pack(side=tk.LEFT, padx=5)
        ttk.Button(ctrl, text="❌ Remove", command=lambda: self.list_remove(self.mword_list)).pack(side=tk.LEFT, padx=5)
        
        tk.Button(f, text="MERGE WORD FILES", bg="#0056b3", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_merge_word).pack(fill=tk.X, padx=40, pady=10)

    def do_merge_word(self):
        if self.mword_list.size() == 0: return
        out = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word", "*.docx")])
        if not out: return
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            new_doc = word.Documents.Add()
            for i in range(self.mword_list.size()):
                file_path = os.path.abspath(self.mword_list.get(i))
                new_doc.Application.Selection.InsertFile(file_path)
                if i < self.mword_list.size() - 1:
                    new_doc.Application.Selection.InsertBreak(7) # wdPageBreak
            new_doc.SaveAs(os.path.abspath(out))
            new_doc.Close()
            word.Quit()
            messagebox.showinfo("Success", "Word files merged to Word document!")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to merge Word files: {e}")

    # 4. CONVERT PDF ↔ WORD
    def init_convert_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="4. PDF↔Word")
        
        # PDF to Word
        f1 = tk.Frame(f, bg="#2d2d2d", bd=2, relief=tk.SUNKEN)
        f1.pack(fill=tk.X, padx=40, pady=15)
        tk.Label(f1, text="PDF to Word", bg="#2d2d2d", fg="white", font=("Segoe UI", 14, "bold")).pack(pady=5)
        self.p2w_file = None
        self.p2w_lbl = tk.Label(f1, text="No PDF selected", bg="#2d2d2d", fg="#007acc")
        self.p2w_lbl.pack()
        btn_p2w = tk.Frame(f1, bg="#2d2d2d")
        btn_p2w.pack()
        ttk.Button(btn_p2w, text="📁 Select PDF", command=lambda: self.select_file("p2w", [("PDF", "*.pdf")])).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_p2w, lambda: self.p2w_file)
        tk.Button(f1, text="CONVERT TO WORD", bg="#17a2b8", fg="white", command=self.do_pdf2word).pack(pady=10)
        
        # Word to PDF
        f2 = tk.Frame(f, bg="#2d2d2d", bd=2, relief=tk.SUNKEN)
        f2.pack(fill=tk.X, padx=40, pady=15)
        tk.Label(f2, text="Word to PDF", bg="#2d2d2d", fg="white", font=("Segoe UI", 14, "bold")).pack(pady=5)
        self.w2p_file = None
        self.w2p_lbl = tk.Label(f2, text="No Word file selected", bg="#2d2d2d", fg="#007acc")
        self.w2p_lbl.pack()
        ttk.Button(f2, text="📁 Select Word", command=lambda: self.select_file("w2p", [("Word", "*.docx *.doc")])).pack(pady=5)
        tk.Button(f2, text="CONVERT TO PDF", bg="#dc3545", fg="white", command=self.do_word2pdf).pack(pady=10)

    def select_file(self, mode, ftypes):
        p = filedialog.askopenfilename(filetypes=ftypes)
        if not p: return
        if mode == "p2w":
            self.p2w_file = p
            self.p2w_lbl.config(text=os.path.basename(p))
        elif mode == "w2p":
            self.w2p_file = p
            self.w2p_lbl.config(text=os.path.basename(p))

    def do_pdf2word(self):
        if not self.p2w_file: return
        out = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word", "*.docx")])
        if not out: return
        try:
            cv = Converter(self.p2w_file)
            cv.convert(out, start=0, end=None)
            cv.close()
            messagebox.showinfo("Success", "PDF converted to Word!")
        except Exception as e:
            messagebox.showerror("Error", str(e))
            
    def do_word2pdf(self):
        if not self.w2p_file: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        try:
            import win32com.client
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(self.w2p_file))
            doc.ExportAsFixedFormat(os.path.abspath(out), 17)
            doc.Close()
            word.Quit()
            messagebox.showinfo("Success", "Word converted to PDF!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # 5. IMAGE ↔ PDF
    def init_image_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="5. Images")
        
        # JPG to PDF
        f1 = tk.Frame(f, bg="#2d2d2d", bd=2, relief=tk.SUNKEN)
        f1.pack(fill=tk.BOTH, expand=True, padx=40, pady=10)
        tk.Label(f1, text="Images to PDF", bg="#2d2d2d", fg="white", font=("Segoe UI", 12, "bold")).pack()
        self.img_list = tk.Listbox(f1, bg="#1e1e1e", fg="white", height=5)
        self.img_list.pack(fill=tk.BOTH, expand=True, padx=20, pady=5)
        
        b1 = tk.Frame(f1, bg="#2d2d2d")
        b1.pack()
        ttk.Button(b1, text="➕ Add Images", command=lambda: self.add_to_list(self.img_list, [("Image", "*.jpg *.png")])).pack(side=tk.LEFT)
        ttk.Button(b1, text="❌ Clear", command=lambda: self.img_list.delete(0, tk.END)).pack(side=tk.LEFT)
        tk.Button(f1, text="CONVERT TO PDF", bg="#e0a800", fg="white", command=self.do_img2pdf).pack(pady=5)
        
        # PDF to JPG
        f2 = tk.Frame(f, bg="#2d2d2d", bd=2, relief=tk.SUNKEN)
        f2.pack(fill=tk.X, padx=40, pady=10)
        tk.Label(f2, text="PDF to Images", bg="#2d2d2d", fg="white", font=("Segoe UI", 12, "bold")).pack()
        self.p2i_file = None
        self.p2i_lbl = tk.Label(f2, text="No PDF", bg="#2d2d2d", fg="#007acc")
        self.p2i_lbl.pack()
        
        btn_f2 = tk.Frame(f2, bg="#2d2d2d")
        btn_f2.pack()
        ttk.Button(btn_f2, text="📁 Select PDF", command=lambda: self.select_p2i()).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f2, lambda: self.p2i_file)
        
        tk.Button(f2, text="EXTRACT PAGES TO IMAGES", bg="#17a2b8", fg="white", command=self.do_pdf2img).pack(pady=5)
        
    def select_p2i(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.p2i_file = p
            self.p2i_lbl.config(text=os.path.basename(p))

    def do_img2pdf(self):
        if self.img_list.size() == 0: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        try:
            imgs = []
            for i in range(self.img_list.size()):
                img = Image.open(self.img_list.get(i))
                if img.mode in ("RGBA", "P"): img = img.convert("RGB")
                imgs.append(img)
            imgs[0].save(out, save_all=True, append_images=imgs[1:])
            messagebox.showinfo("Success", "Images to PDF complete!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def do_pdf2img(self):
        if not self.p2i_file: return
        out_dir = filedialog.askdirectory(title="Select Output Folder")
        if not out_dir: return
        try:
            doc = fitz.open(self.p2i_file)
            base = os.path.splitext(os.path.basename(self.p2i_file))[0]
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                pix.save(os.path.join(out_dir, f"{base}_page_{i+1}.jpg"))
            messagebox.showinfo("Success", f"Saved {len(doc)} images!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # 6. PDF TO PPTX
    def init_pptx_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="6. To PPTX")
        tk.Label(f, text="Convert PDF pages to PowerPoint slides", bg="#1e1e1e", fg="white", font=("Segoe UI", 12)).pack(pady=20)
        self.pptx_file = None
        self.pptx_lbl = tk.Label(f, text="No PDF", bg="#1e1e1e", fg="#007acc")
        self.pptx_lbl.pack()
        
        btn_f = tk.Frame(f, bg="#1e1e1e")
        btn_f.pack(pady=5)
        ttk.Button(btn_f, text="📁 Select PDF", command=lambda: self.select_pptx()).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f, lambda: self.pptx_file)
        
        tk.Button(f, text="CONVERT TO PPTX", bg="#fd7e14", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_pdf2pptx).pack(fill=tk.X, padx=100, pady=20)
        
    def select_pptx(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.pptx_file = p
            self.pptx_lbl.config(text=os.path.basename(p))
            
    def do_pdf2pptx(self):
        if not self.pptx_file: return
        out = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if not out: return
        try:
            doc = fitz.open(self.pptx_file)
            prs = Presentation()
            blank = prs.slide_layouts[6]
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=fitz.Matrix(2,2))
                temp = f"temp_slide_{i}.png"
                pix.save(temp)
                slide = prs.slides.add_slide(blank)
                slide.shapes.add_picture(temp, 0, 0, width=prs.slide_width, height=prs.slide_height)
                os.remove(temp)
            prs.save(out)
            messagebox.showinfo("Success", "PowerPoint created!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # 7. SPLIT
    def init_split_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="7. Split")
        self.split_file = None
        self.split_lbl = tk.Label(f, text="No PDF", bg="#1e1e1e", fg="#007acc")
        self.split_lbl.pack(pady=10)
        
        btn_f = tk.Frame(f, bg="#1e1e1e")
        btn_f.pack(pady=5)
        ttk.Button(btn_f, text="📁 Select PDF", command=self.sel_split).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f, lambda: self.split_file)
        
        tk.Label(f, text="Split Range (e.g. 1 to 5):", bg="#1e1e1e", fg="white").pack(pady=5)
        rf = tk.Frame(f, bg="#1e1e1e")
        rf.pack()
        self.sp_start = ttk.Entry(rf, width=5)
        self.sp_start.pack(side=tk.LEFT)
        tk.Label(rf, text=" to ", bg="#1e1e1e", fg="white").pack(side=tk.LEFT)
        self.sp_end = ttk.Entry(rf, width=5)
        self.sp_end.pack(side=tk.LEFT)
        
        tk.Button(f, text="EXTRACT RANGE", bg="#6f42c1", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_split_range).pack(fill=tk.X, padx=100, pady=10)
        tk.Label(f, text="OR", bg="#1e1e1e", fg="gray").pack()
        tk.Button(f, text="SPLIT ALL PAGES TO SINGLE FILES", bg="#20c997", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_split_all).pack(fill=tk.X, padx=100, pady=10)

    def sel_split(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.split_file = p
            self.split_lbl.config(text=os.path.basename(p))

    def do_split_range(self):
        if not self.split_file: return
        try:
            st = int(self.sp_start.get()) - 1
            ed = int(self.sp_end.get()) - 1
            out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
            if not out: return
            
            reader = PdfReader(self.split_file)
            writer = PdfWriter()
            for i in range(st, ed + 1):
                writer.add_page(reader.pages[i])
            with open(out, "wb") as f:
                writer.write(f)
            messagebox.showinfo("Success", "Range extracted!")
        except Exception as e:
            messagebox.showerror("Error", f"Invalid range or error:\n{e}")

    def do_split_all(self):
        if not self.split_file: return
        out_dir = filedialog.askdirectory()
        if not out_dir: return
        try:
            reader = PdfReader(self.split_file)
            base = os.path.splitext(os.path.basename(self.split_file))[0]
            for i in range(len(reader.pages)):
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                with open(os.path.join(out_dir, f"{base}_page_{i+1}.pdf"), "wb") as f:
                    writer.write(f)
            messagebox.showinfo("Success", "All pages split!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # 8. LOCK / UNLOCK
    def init_lock_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="8. Lock/Unlock")
        self.lu_file = None
        self.lu_lbl = tk.Label(f, text="No PDF", bg="#1e1e1e", fg="#007acc")
        self.lu_lbl.pack(pady=10)
        
        btn_f = tk.Frame(f, bg="#1e1e1e")
        btn_f.pack(pady=5)
        ttk.Button(btn_f, text="📁 Select PDF", command=self.sel_lu).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f, lambda: self.lu_file)
        
        pf = tk.Frame(f, bg="#1e1e1e")
        pf.pack(pady=10)
        tk.Label(pf, text="Password:", bg="#1e1e1e", fg="white").pack(side=tk.LEFT)
        self.lu_pass = ttk.Entry(pf, show="*")
        self.lu_pass.pack(side=tk.LEFT, padx=5)
        
        tk.Button(f, text="🔒 LOCK PDF", bg="#dc3545", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_lock).pack(fill=tk.X, padx=100, pady=5)
        tk.Button(f, text="🔓 UNLOCK PDF", bg="#28a745", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_unlock).pack(fill=tk.X, padx=100, pady=5)
        
        tk.Label(f, text="--- Rotate ---", bg="#1e1e1e", fg="gray", font=("Segoe UI", 12)).pack(pady=10)
        rf = tk.Frame(f, bg="#1e1e1e")
        rf.pack()
        self.rot_var = tk.IntVar(value=90)
        tk.Radiobutton(rf, text="90°", variable=self.rot_var, value=90, bg="#1e1e1e", fg="white", selectcolor="#1e1e1e").pack(side=tk.LEFT)
        tk.Radiobutton(rf, text="180°", variable=self.rot_var, value=180, bg="#1e1e1e", fg="white", selectcolor="#1e1e1e").pack(side=tk.LEFT)
        tk.Radiobutton(rf, text="270°", variable=self.rot_var, value=270, bg="#1e1e1e", fg="white", selectcolor="#1e1e1e").pack(side=tk.LEFT)
        tk.Button(f, text="ROTATE", bg="#17a2b8", fg="white", command=self.do_rotate).pack(pady=5)

    def sel_lu(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.lu_file = p
            self.lu_lbl.config(text=os.path.basename(p))

    def do_lock(self):
        if not self.lu_file or not self.lu_pass.get(): return
        out = filedialog.asksaveasfilename(defaultextension=".pdf")
        if not out: return
        try:
            r = PdfReader(self.lu_file)
            w = PdfWriter()
            for p in r.pages: w.add_page(p)
            w.encrypt(self.lu_pass.get())
            with open(out, "wb") as f: w.write(f)
            messagebox.showinfo("Success", "Locked!")
        except Exception as e: messagebox.showerror("Error", str(e))
        
    def do_unlock(self):
        if not self.lu_file or not self.lu_pass.get(): return
        out = filedialog.asksaveasfilename(defaultextension=".pdf")
        if not out: return
        try:
            r = PdfReader(self.lu_file)
            r.decrypt(self.lu_pass.get())
            w = PdfWriter()
            for p in r.pages: w.add_page(p)
            with open(out, "wb") as f: w.write(f)
            messagebox.showinfo("Success", "Unlocked!")
        except Exception as e: messagebox.showerror("Error", "Failed to unlock. Incorrect password or error.")
        
    def do_rotate(self):
        if not self.lu_file: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf")
        if not out: return
        try:
            r = PdfReader(self.lu_file)
            w = PdfWriter()
            for p in r.pages:
                p.rotate(self.rot_var.get())
                w.add_page(p)
            with open(out, "wb") as f: w.write(f)
            messagebox.showinfo("Success", "Rotated!")
        except Exception as e: messagebox.showerror("Error", str(e))

    # 9. ANNOTATE (Highlight text basic alternative) & Scan
    def init_annotate_tab(self):
        f = ttk.Frame(self.notebook)
        self.notebook.add(f, text="9. Annotate")
        tk.Label(f, text="Extract text or Highlight a specific word in PDF", bg="#1e1e1e", fg="white").pack(pady=10)
        
        self.ano_file = None
        self.ano_lbl = tk.Label(f, text="No PDF", bg="#1e1e1e", fg="#007acc")
        self.ano_lbl.pack()
        
        btn_f = tk.Frame(f, bg="#1e1e1e")
        btn_f.pack(pady=5)
        ttk.Button(btn_f, text="📁 Select PDF", command=self.sel_ano).pack(side=tk.LEFT, padx=5)
        self.create_preview_btn(btn_f, lambda: self.ano_file)
        
        tk.Button(f, text="EXTRACT ALL TEXT TO .TXT", bg="#6f42c1", fg="white", font=("Segoe UI", 12, "bold"), command=self.do_extract).pack(fill=tk.X, padx=100, pady=10)
        
        hf = tk.Frame(f, bg="#1e1e1e", bd=2, relief=tk.SUNKEN)
        hf.pack(fill=tk.X, padx=40, pady=10)
        tk.Label(hf, text="Search & Highlight Word:", bg="#1e1e1e", fg="white").pack(pady=5)
        self.hl_word = ttk.Entry(hf)
        self.hl_word.pack(pady=5)
        tk.Button(hf, text="HIGHLIGHT & SAVE PDF", bg="#e83e8c", fg="white", command=self.do_highlight).pack(pady=10)
        
    def sel_ano(self):
        p = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if p:
            self.ano_file = p
            self.ano_lbl.config(text=os.path.basename(p))
            
    def do_extract(self):
        if not self.ano_file: return
        out = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text", "*.txt")])
        if not out: return
        try:
            doc = fitz.open(self.ano_file)
            text = ""
            for p in doc: text += p.get_text() + "\n"
            with open(out, "w", encoding="utf-8") as f: f.write(text)
            messagebox.showinfo("Success", "Text extracted!")
        except Exception as e: messagebox.showerror("Error", str(e))
        
    def do_highlight(self):
        if not self.ano_file or not self.hl_word.get(): return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")])
        if not out: return
        try:
            doc = fitz.open(self.ano_file)
            for page in doc:
                rl = page.search_for(self.hl_word.get())
                for inst in rl:
                    page.add_highlight_annot(inst)
            doc.save(out)
            messagebox.showinfo("Success", f"Highlighted all instances of '{self.hl_word.get()}'!")
        except Exception as e:
            messagebox.showerror("Error", str(e))

if __name__ == "__main__":
    root = tk.Tk()
    app = UltimatePDFToolkit(root)
    root.mainloop()
